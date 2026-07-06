import hashlib
import json
import re
from pathlib import Path

import faiss
import numpy as np
from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader
from sqlalchemy.orm import Session

from .config import INDEX_DIR, settings
from .models import Document, DocumentChunk


FALLBACK = "I could not find relevant information in faculty uploaded study materials."
_model = None


def extract_pages(path: Path, suffix: str) -> list[tuple[int, str]]:
    if suffix == ".pdf":
        return [(i + 1, page.extract_text() or "") for i, page in enumerate(PdfReader(str(path)).pages)]
    if suffix == ".docx":
        doc = DocxDocument(str(path))
        return [(1, "\n".join(p.text for p in doc.paragraphs))]
    if suffix == ".pptx":
        presentation = Presentation(str(path))
        pages = []
        for i, slide in enumerate(presentation.slides):
            text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
            pages.append((i + 1, text))
        return pages
    raise ValueError("Unsupported file type")


def chunk_pages(pages: list[tuple[int, str]], size: int = 900, overlap: int = 150) -> list[tuple[int, str]]:
    chunks = []
    for page, text in pages:
        clean = re.sub(r"\s+", " ", text).strip()
        start = 0
        while start < len(clean):
            chunk = clean[start:start + size].strip()
            if chunk:
                chunks.append((page, chunk))
            start += size - overlap
    return chunks


def _hash_embeddings(texts: list[str], dims: int = 384) -> np.ndarray:
    vectors = np.zeros((len(texts), dims), dtype="float32")
    for row, text in enumerate(texts):
        for token in re.findall(r"[a-z0-9]+", text.lower()):
            digest = hashlib.blake2b(token.encode(), digest_size=8).digest()
            value = int.from_bytes(digest, "little")
            vectors[row, value % dims] += 1.0 if value & 1 else -1.0
    faiss.normalize_L2(vectors)
    return vectors


def embed(texts: list[str]) -> np.ndarray:
    global _model
    if settings.embedding_provider.lower() == "sentence-transformers":
        try:
            if _model is None:
                from sentence_transformers import SentenceTransformer
                _model = SentenceTransformer(settings.embedding_model)
            return np.asarray(_model.encode(texts, normalize_embeddings=True), dtype="float32")
        except Exception:
            return _hash_embeddings(texts)
    return _hash_embeddings(texts)


def rebuild_subject_index(db: Session, subject_id: int) -> int:
    rows = (
        db.query(DocumentChunk)
        .join(Document)
        .filter(Document.subject_id == subject_id)
        .order_by(DocumentChunk.id)
        .all()
    )
    index_path = INDEX_DIR / f"subject_{subject_id}.faiss"
    meta_path = INDEX_DIR / f"subject_{subject_id}.json"
    if not rows:
        index_path.unlink(missing_ok=True)
        meta_path.unlink(missing_ok=True)
        return 0
    vectors = embed([row.content for row in rows])
    index = faiss.IndexFlatIP(vectors.shape[1])
    index.add(vectors)
    faiss.write_index(index, str(index_path))
    meta_path.write_text(json.dumps([row.id for row in rows]), encoding="utf-8")
    return len(rows)


def retrieve(db: Session, subject_id: int, question: str, limit: int = 5) -> list[tuple[DocumentChunk, float]]:
    index_path = INDEX_DIR / f"subject_{subject_id}.faiss"
    meta_path = INDEX_DIR / f"subject_{subject_id}.json"
    if not index_path.exists() or not meta_path.exists():
        rebuild_subject_index(db, subject_id)
    if not index_path.exists():
        return []
    index = faiss.read_index(str(index_path))
    ids = json.loads(meta_path.read_text(encoding="utf-8"))
    scores, positions = index.search(embed([question]), min(limit, len(ids)))
    results = []
    for score, pos in zip(scores[0], positions[0]):
        if pos >= 0:
            chunk = db.get(DocumentChunk, ids[pos])
            if chunk:
                results.append((chunk, float(score)))
    return results


def generate(prompt: str) -> str | None:
    if not settings.gemini_api_key:
        return None
    try:
        from google import genai
        client = genai.Client(api_key=settings.gemini_api_key)
        response = client.models.generate_content(model=settings.gemini_model, contents=prompt)
        return response.text
    except Exception:
        return None


def answer_question(db: Session, subject_id: int, question: str) -> tuple[str, DocumentChunk | None]:
    results = retrieve(db, subject_id, question)
    if not results or results[0][1] < 0.12:
        return FALLBACK, None
    context = "\n\n".join(
        f"[Source: {chunk.document.filename}, page {chunk.page_number}]\n{chunk.content}"
        for chunk, _ in results
    )
    prompt = (
        "Answer only from the supplied study material. If the answer is absent, reply exactly: "
        f"{FALLBACK}\n\nQuestion: {question}\n\nMaterial:\n{context}"
    )
    answer = generate(prompt)
    if not answer:
        answer = results[0][0].content
    return answer.strip(), results[0][0]


def subject_text(db: Session, subject_id: int, max_chars: int = 50000) -> str:
    rows = (
        db.query(DocumentChunk)
        .join(Document)
        .filter(Document.subject_id == subject_id)
        .order_by(DocumentChunk.id)
        .all()
    )
    return "\n\n".join(row.content for row in rows)[:max_chars]


def make_quiz(db: Session, subject_id: int, difficulty: str, count: int) -> list[dict]:
    material = subject_text(db, subject_id)
    if not material:
        raise ValueError("No study materials found for this subject")
    prompt = (
        f"Create exactly {count} {difficulty} multiple-choice questions from this material. "
        'Return only JSON array with keys "question", "options" (4 strings), "answer", "explanation".\n'
        + material
    )
    generated = generate(prompt)
    if generated:
        try:
            return json.loads(re.sub(r"^```json|```$", "", generated.strip(), flags=re.MULTILINE).strip())
        except json.JSONDecodeError:
            pass
    sentences = [s.strip() for s in re.split(r"(?<=[.!?])\s+", material) if len(s.split()) > 7]
    if not sentences:
        sentences = [material[:250]]
    quiz = []
    for i in range(count):
        sentence = sentences[i % len(sentences)]
        quiz.append({
            "question": f"Which statement is supported by the uploaded material? ({i + 1})",
            "options": [sentence, "None of the uploaded material discusses this.", "The opposite statement is true.", "This cannot be determined."],
            "answer": sentence,
            "explanation": "This statement appears directly in the uploaded study material.",
        })
    return quiz


def make_summary(db: Session, subject_id: int) -> str:
    material = subject_text(db, subject_id)
    if not material:
        raise ValueError("No study materials found for this subject")
    generated = generate("Summarize these study materials into clear headings and bullet points:\n" + material)
    if generated:
        return generated.strip()
    sentences = [s.strip() for s in re.split(r"(?<=[.!?])\s+", material) if len(s.strip()) > 30]
    return "\n".join(f"- {sentence}" for sentence in sentences[:12])
